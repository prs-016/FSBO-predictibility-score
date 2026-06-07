"""Build the FBSO Predictor pitch deck (lightning, ~5 min, coach-facing)."""
from __future__ import annotations

import math
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# ─── Theme ────────────────────────────────────────────────────────────────
DARK    = RGBColor(0x0d, 0x0f, 0x14)   # background
PANEL   = RGBColor(0x16, 0x19, 0x20)   # raised panels
PANEL_2 = RGBColor(0x1e, 0x22, 0x29)   # deeper panel
TEXT    = RGBColor(0xe8, 0xec, 0xf1)
DIM     = RGBColor(0x88, 0x92, 0xa2)
ACCENT  = RGBColor(0x4c, 0xc9, 0xf0)   # cyan glow
PINK    = RGBColor(0xff, 0x6e, 0xc7)   # brain-image magenta
GOLD    = RGBColor(0xff, 0xcd, 0x00)   # UCSD gold
GOOD    = RGBColor(0x2e, 0xcc, 0x71)

FONT_DISPLAY = "Helvetica"
FONT_BODY = "Helvetica"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def make_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def fill_background(slide, color: RGBColor) -> None:
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    spTree = rect._element.getparent()
    spTree.remove(rect._element)
    spTree.insert(2, rect._element)


def add_text(
    slide, x, y, w, h, text: str,
    *, size: int = 18, color: RGBColor = TEXT, bold: bool = False,
    font: str = FONT_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
    italic: bool = False,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_accent_bar(slide, x, y, w, h, color: RGBColor = ACCENT) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False


def add_glow_circle(slide, cx, cy, r, color: RGBColor, alpha_pct: int = 100) -> None:
    """Cyan/pink glowing circle — for the brain/neuron motif."""
    diameter = r * 2
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, diameter, diameter)
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.fill.background()
    c.shadow.inherit = False


def add_neuron_motif(slide, anchor_x_emu: int, anchor_y_emu: int, scale: float = 1.0) -> None:
    """Decorative neuron-cluster motif — small dots connected by thin lines,
    evoking the neural-network/brain aesthetic from the source image.
    """
    base = int(Inches(0.18) * scale)
    # 5 nodes in a soft cluster
    nodes = [
        (anchor_x_emu + int(Inches(0.0) * scale),  anchor_y_emu + int(Inches(0.0) * scale),  base, ACCENT),
        (anchor_x_emu + int(Inches(0.9) * scale),  anchor_y_emu - int(Inches(0.5) * scale),  base // 2, PINK),
        (anchor_x_emu + int(Inches(1.4) * scale),  anchor_y_emu + int(Inches(0.2) * scale),  base // 2, ACCENT),
        (anchor_x_emu + int(Inches(0.6) * scale),  anchor_y_emu + int(Inches(0.9) * scale),  base // 3, PINK),
        (anchor_x_emu + int(Inches(1.7) * scale),  anchor_y_emu + int(Inches(1.0) * scale),  base // 3, ACCENT),
    ]
    # Connector lines
    for i, (x1, y1, _, _) in enumerate(nodes):
        for x2, y2, _, _ in nodes[i + 1: i + 3]:
            ln = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = STRAIGHT
            ln.line.color.rgb = ACCENT
            ln.line.width = Emu(int(Inches(0.012) * scale))
            ln.shadow.inherit = False
    for x, y, r, col in nodes:
        add_glow_circle(slide, x, y, r, col)


def add_footer(slide, slide_num: int, total: int, label: str = "FBSO Predictor · UCSD Triton Analytics") -> None:
    add_text(
        slide, Inches(0.5), SLIDE_H - Inches(0.4), Inches(8), Inches(0.3),
        label, size=10, color=DIM, italic=True,
    )
    add_text(
        slide, SLIDE_W - Inches(1.5), SLIDE_H - Inches(0.4), Inches(1), Inches(0.3),
        f"{slide_num} / {total}", size=10, color=DIM, align=PP_ALIGN.RIGHT,
    )


def add_notes(slide, notes: str) -> None:
    slide.notes_slide.notes_text_frame.text = notes


# ─── Slide builders ───────────────────────────────────────────────────────


def slide_title(prs: Presentation, total: int) -> None:
    s = add_blank(prs)
    fill_background(s, DARK)
    # Decorative motif top-right
    add_neuron_motif(s, int(Inches(10.0)), int(Inches(1.5)), scale=1.4)
    # Subtle accent column on left
    add_accent_bar(s, Inches(0), Inches(0), Inches(0.18), SLIDE_H)
    # Title block
    add_text(
        s, Inches(0.9), Inches(2.4), Inches(11), Inches(1.6),
        "FBSO Predictor", size=70, color=ACCENT, bold=True, font=FONT_DISPLAY,
    )
    add_text(
        s, Inches(0.9), Inches(3.7), Inches(11), Inches(1.0),
        "AI that reads their setter — before the swing.",
        size=28, color=TEXT, font=FONT_BODY,
    )
    add_text(
        s, Inches(0.9), Inches(4.8), Inches(11), Inches(0.5),
        "First-Ball Side-Out scouting, live on the bench.",
        size=18, color=DIM, italic=True,
    )
    add_text(
        s, Inches(0.9), Inches(6.5), Inches(11), Inches(0.4),
        "UC San Diego Volleyball  ·  Triton Analytics",
        size=14, color=GOLD, bold=True,
    )
    add_footer(s, 1, total)
    add_notes(
        s,
        "Open strong. Single sentence: 'We built a tool that predicts what the opposing setter will call — "
        "before the ball is swung — based on the same patterns your scouts already watch for.' "
        "Hold the title for 5–8 seconds. Don't over-explain yet.",
    )


def slide_problem(prs: Presentation, total: int) -> None:
    s = add_blank(prs)
    fill_background(s, DARK)
    add_accent_bar(s, Inches(0), Inches(0), Inches(0.18), SLIDE_H)

    add_text(
        s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.7),
        "The 20-second problem", size=40, color=ACCENT, bold=True, font=FONT_DISPLAY,
    )
    add_text(
        s, Inches(0.6), Inches(1.5), Inches(12), Inches(0.5),
        "Between every rally, the coach has one job:",
        size=20, color=DIM,
    )
    add_text(
        s, Inches(0.6), Inches(2.1), Inches(12), Inches(0.7),
        "Call the right defense for the next swing.",
        size=28, color=TEXT, bold=True,
    )

    # Three-card layout
    card_y = Inches(3.4)
    card_w = Inches(3.95)
    card_h = Inches(3.0)
    gap = Inches(0.2)
    start_x = Inches(0.6)

    cards = [
        ("TODAY", "Film study + memory.",
         "Pre-match notes. A guess on every rally. Easy to read the obvious setter; hard to read late-set or out-of-system."),
        ("THE GAP", "15–30 seconds per rally.",
         "Not enough time to look up tendencies by rotation, score, recent pattern — even when you know exactly what to ask."),
        ("THE OPPORTUNITY", "Predictable patterns.",
         "Setters lean on their best hitters under pressure. Streaks, score state, rotation, and recent calls all shift the odds — measurably."),
    ]
    for i, (tag, headline, body) in enumerate(cards):
        x = start_x + (card_w + gap) * i
        # card background
        rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, card_y, card_w, card_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = PANEL
        rect.line.color.rgb = PANEL_2
        rect.line.width = Emu(int(Inches(0.01)))
        rect.shadow.inherit = False
        # tag
        add_text(s, x + Inches(0.4), card_y + Inches(0.35), card_w - Inches(0.8), Inches(0.4),
                 tag, size=11, color=DIM, bold=True)
        # headline
        add_text(s, x + Inches(0.4), card_y + Inches(0.8), card_w - Inches(0.8), Inches(0.9),
                 headline, size=22, color=TEXT, bold=True)
        # body
        add_text(s, x + Inches(0.4), card_y + Inches(1.85), card_w - Inches(0.8), card_h - Inches(2.0),
                 body, size=13, color=DIM)

    add_footer(s, 2, total)
    add_notes(
        s,
        "Frame the 'why now'. The window between rallies is short. Even the best coach can't simultaneously "
        "read score, rotation, streak, AND recall a film-study note. We close that gap with the same kind of "
        "pattern recognition the brain does — just faster and over more data.",
    )


def slide_solution(prs: Presentation, total: int) -> None:
    s = add_blank(prs)
    fill_background(s, DARK)
    add_accent_bar(s, Inches(0), Inches(0), Inches(0.18), SLIDE_H)
    add_neuron_motif(s, int(Inches(11.5)), int(Inches(0.8)), scale=0.9)

    add_text(
        s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.7),
        "What it does", size=40, color=ACCENT, bold=True, font=FONT_DISPLAY,
    )
    add_text(
        s, Inches(0.6), Inches(1.5), Inches(12), Inches(0.5),
        "Bench-side, one second after their pass.",
        size=20, color=DIM,
    )

    # Flow boxes: reception → predict → call
    flow_y = Inches(2.6)
    flow_h = Inches(2.0)
    bw = Inches(3.6)
    gap = Inches(0.45)
    fx = [Inches(0.6), Inches(0.6) + bw + gap, Inches(0.6) + (bw + gap) * 2]
    flow = [
        ("THEIR RECEPTION", "Pass + setter ready", PANEL_2),
        ("OUR MODEL", "13 features → top-K probabilities", ACCENT),
        ("YOUR CALL", "Block alignment + defensive set", GOLD),
    ]
    for i, ((tag, body, color)) in enumerate(zip(["THEIR RECEPTION", "OUR MODEL", "YOUR CALL"],
                                                  ["Pass + setter ready",
                                                   "13 features → top-K probabilities",
                                                   "Block alignment + defensive set"],
                                                  [PANEL_2, ACCENT, GOLD])):
        rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, fx[i], flow_y, bw, flow_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = PANEL
        rect.line.color.rgb = color
        rect.line.width = Emu(int(Inches(0.025)))
        rect.shadow.inherit = False
        add_text(s, fx[i] + Inches(0.35), flow_y + Inches(0.4), bw - Inches(0.7), Inches(0.4),
                 tag, size=11, color=color, bold=True)
        add_text(s, fx[i] + Inches(0.35), flow_y + Inches(0.95), bw - Inches(0.7), Inches(1.0),
                 body, size=18, color=TEXT, bold=True)

    # Arrows between
    for i in range(2):
        ax = fx[i] + bw + Inches(0.02)
        ay = flow_y + Inches(0.9)
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, ay, gap - Inches(0.04), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT
        arrow.line.fill.background()
        arrow.shadow.inherit = False

    # Sample output preview
    sample_y = Inches(5.05)
    add_text(s, Inches(0.6), sample_y, Inches(12), Inches(0.4),
             "WHAT THE BENCH SEES", size=11, color=DIM, bold=True)
    rows = [("Back",   "62%", ACCENT, 0.62),
            ("Front",  "32%", DIM,    0.32),
            ("Middle", "6%",  DIM,    0.06)]
    bar_x = Inches(0.6)
    bar_y0 = Inches(5.5)
    bar_w_max = Inches(8.0)
    row_h = Inches(0.45)
    for i, (cat, prob, color, frac) in enumerate(rows):
        y = bar_y0 + i * (row_h + Inches(0.1))
        # background
        bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bar_x, y, bar_w_max, row_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = PANEL_2
        bg.line.fill.background()
        bg.shadow.inherit = False
        # filled portion
        fill_w = int(bar_w_max * frac)
        fg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bar_x, y, fill_w, row_h)
        fg.fill.solid()
        fg.fill.fore_color.rgb = color
        fg.line.fill.background()
        fg.shadow.inherit = False
        # category label inside bar
        add_text(s, bar_x + Inches(0.25), y, bar_w_max, row_h,
                 cat, size=15, color=TEXT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        # probability on right
        add_text(s, bar_x + bar_w_max + Inches(0.2), y, Inches(1.2), row_h,
                 prob, size=15, color=color, bold=True, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(10.0), bar_y0, Inches(3), Inches(0.4),
             "ONE-SECOND LATENCY", size=11, color=DIM, bold=True)
    add_text(s, Inches(10.0), bar_y0 + Inches(0.5), Inches(3), Inches(1.5),
             "Live, on a courtside laptop. No cloud round-trip.",
             size=14, color=TEXT)

    add_footer(s, 3, total)
    add_notes(
        s,
        "Walk the flow left to right. Emphasize: 'You don't change anything about how you scout — we read what "
        "the scout is already typing into DataVolley. The model adds a second voice in the room.'",
    )


def slide_credibility(prs: Presentation, total: int) -> None:
    s = add_blank(prs)
    fill_background(s, DARK)
    add_accent_bar(s, Inches(0), Inches(0), Inches(0.18), SLIDE_H)
    add_neuron_motif(s, int(Inches(11.0)), int(Inches(0.9)), scale=1.0)

    add_text(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.7),
             "Built on real opponents — not theory",
             size=36, color=ACCENT, bold=True, font=FONT_DISPLAY)

    # Big number column
    add_text(s, Inches(0.6), Inches(2.0), Inches(3.5), Inches(2.5),
             "33", size=170, color=PINK, bold=True, font=FONT_DISPLAY)
    add_text(s, Inches(0.6), Inches(4.4), Inches(3.5), Inches(0.5),
             "opponent-specific models", size=18, color=TEXT, bold=True)
    add_text(s, Inches(0.6), Inches(5.0), Inches(3.5), Inches(1.5),
             "Each team has its own model trained on its own attacks. "
             "What they call after a perfect pass is different from what they call after a scramble.",
             size=12, color=DIM)

    # Right panel: opponent list
    rx = Inches(4.6)
    rw = Inches(8.2)
    rh = Inches(4.7)
    ry = Inches(1.7)
    rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rx, ry, rw, rh)
    rect.fill.solid()
    rect.fill.fore_color.rgb = PANEL
    rect.line.color.rgb = PANEL_2
    rect.line.width = Emu(int(Inches(0.01)))
    rect.shadow.inherit = False

    add_text(s, rx + Inches(0.4), ry + Inches(0.3), rw - Inches(0.8), Inches(0.4),
             "BIG WEST + WCC OPPONENTS COVERED", size=11, color=ACCENT, bold=True)

    opponents = [
        "Cal Poly", "Long Beach State", "Cal State Northridge",
        "UC Davis", "UC Irvine", "UC Riverside",
        "UC Santa Barbara", "Hawaii", "CSU Bakersfield",
        "Stanford", "UC Berkeley", "UCLA",
        "Pepperdine", "U. San Diego", "U. Pacific",
        "+ 18 more teams",
    ]
    cols = 4
    col_w = (rw - Inches(0.8)) / cols
    row_h = Inches(0.42)
    for i, name in enumerate(opponents):
        col = i % cols
        row = i // cols
        x = rx + Inches(0.4) + col_w * col
        y = ry + Inches(0.9) + row_h * row
        bullet = "›" if name.startswith("+") else "▸"
        col_color = GOLD if name.startswith("+") else TEXT
        add_text(s, x, y, col_w - Inches(0.1), row_h,
                 f"{bullet}  {name}", size=13,
                 color=col_color,
                 bold=name.startswith("+"))

    add_text(s, rx + Inches(0.4), ry + rh - Inches(0.85), rw - Inches(0.8), Inches(0.4),
             "MODEL", size=11, color=ACCENT, bold=True)
    add_text(s, rx + Inches(0.4), ry + rh - Inches(0.45), rw - Inches(0.8), Inches(0.4),
             "BiLSTM + Gradient Boosting ensemble  ·  4 attack zones: Front / Middle / Back / Pipe",
             size=13, color=DIM)

    add_footer(s, 4, total)
    add_notes(
        s,
        "Credibility slide. The 'big 33' is the answer to 'but who?'. We have a per-opponent model for nearly "
        "every team you'll face this season. Each one learns that team's specific tendencies — not a one-size-fits-all "
        "model. The ensemble (BiLSTM + GBM) is the same kind of stack used by pro analytics teams.",
    )


def slide_workflow(prs: Presentation, total: int) -> None:
    s = add_blank(prs)
    fill_background(s, DARK)
    add_accent_bar(s, Inches(0), Inches(0), Inches(0.18), SLIDE_H)

    add_text(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.7),
             "Bench-ready today", size=40, color=ACCENT, bold=True, font=FONT_DISPLAY)
    add_text(s, Inches(0.6), Inches(1.5), Inches(12), Inches(0.5),
             "Two ways to drive the prediction — both available now.",
             size=20, color=DIM)

    # Two large mode cards
    card_y = Inches(2.5)
    card_w = Inches(5.95)
    card_h = Inches(4.0)
    gap = Inches(0.3)

    modes = [
        (
            "LIVE",
            "Reads your scout in real time.",
            "Tail the DataVolley file your stats team already writes. Each "
            "opponent reception fires a fresh prediction in under one second. "
            "Coach watches the bar chart move; never types anything.",
            ACCENT,
        ),
        (
            "MANUAL",
            "Type the situation. Get the read.",
            "No scout running? Coach or assistant clicks dropdowns: score, "
            "rotation, last 5 attacks. Predict. Use it in timeouts, pre-match "
            "prep, or as a teaching tool.",
            PINK,
        ),
    ]

    for i, (tag, headline, body, color) in enumerate(modes):
        x = Inches(0.6) + (card_w + gap) * i
        rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, card_y, card_w, card_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = PANEL
        rect.line.color.rgb = color
        rect.line.width = Emu(int(Inches(0.03)))
        rect.shadow.inherit = False
        add_text(s, x + Inches(0.5), card_y + Inches(0.4), card_w - Inches(1.0), Inches(0.4),
                 tag, size=12, color=color, bold=True)
        add_text(s, x + Inches(0.5), card_y + Inches(0.9), card_w - Inches(1.0), Inches(1.2),
                 headline, size=26, color=TEXT, bold=True)
        add_text(s, x + Inches(0.5), card_y + Inches(2.2), card_w - Inches(1.0), card_h - Inches(2.4),
                 body, size=15, color=DIM)

    # Three icons row
    icons_y = Inches(6.7)
    icons = [
        ("⚡  < 1 sec / prediction", 0),
        ("💻  Runs on the bench laptop", 4.5),
        ("🔒  No video upload. Local only.", 9.0),
    ]
    for label, dx in icons:
        add_text(s, Inches(0.6 + dx), icons_y, Inches(4.5), Inches(0.4),
                 label, size=14, color=GOOD, bold=True)

    add_footer(s, 5, total)
    add_notes(
        s,
        "Address the 'how does this fit our existing process' worry. Live mode = zero workflow change. "
        "Manual mode = anyone can use it without a scout. Latency point matters: it's local, not in the cloud — "
        "no internet dependency in the gym.",
    )


def slide_ask(prs: Presentation, total: int) -> None:
    s = add_blank(prs)
    fill_background(s, DARK)
    add_accent_bar(s, Inches(0), Inches(0), Inches(0.18), SLIDE_H)
    add_neuron_motif(s, int(Inches(11.4)), int(Inches(1.0)), scale=0.9)

    add_text(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.7),
             "What we need from you",
             size=40, color=ACCENT, bold=True, font=FONT_DISPLAY)
    add_text(s, Inches(0.6), Inches(1.5), Inches(12), Inches(0.5),
             "To pilot this in the spring season.",
             size=20, color=DIM)

    asks = [
        ("01",
         "One pilot match.",
         "Pick a Big West opponent we already have a model for. Park "
         "the laptop next to the scout. Use it on side-out defense calls. "
         "Compare your bench notes after the match."),
        ("02",
         "Scout team access.",
         "Five minutes of your DataVolley operator's time to point us at "
         "their safety-file location. Zero workflow change for them."),
        ("03",
         "Athletic dept sign-off.",
         "A green light to use it during competition. We've designed it to "
         "be a passive advisor — coach decisions stay coach decisions."),
    ]

    card_y = Inches(2.4)
    card_w = Inches(3.95)
    card_h = Inches(4.0)
    gap = Inches(0.2)
    start_x = Inches(0.6)
    for i, (num, headline, body) in enumerate(asks):
        x = start_x + (card_w + gap) * i
        rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, card_y, card_w, card_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = PANEL
        rect.line.color.rgb = ACCENT
        rect.line.width = Emu(int(Inches(0.02)))
        rect.shadow.inherit = False
        add_text(s, x + Inches(0.4), card_y + Inches(0.35), card_w - Inches(0.8), Inches(1.2),
                 num, size=42, color=PINK, bold=True, font=FONT_DISPLAY)
        add_text(s, x + Inches(0.4), card_y + Inches(1.55), card_w - Inches(0.8), Inches(1.0),
                 headline, size=22, color=TEXT, bold=True)
        add_text(s, x + Inches(0.4), card_y + Inches(2.7), card_w - Inches(0.8), card_h - Inches(2.9),
                 body, size=13, color=DIM)

    # Pre-empt: "what about..."
    add_text(s, Inches(0.6), Inches(6.65), Inches(12), Inches(0.5),
             "Designed to advise — not replace. No player data leaves the laptop. No video required.",
             size=13, color=DIM, italic=True)

    add_footer(s, 6, total)
    add_notes(
        s,
        "Three concrete asks. Make it small: one match, five minutes, one signature. Address the two real "
        "objections coaches will have: (a) 'I don't want a computer telling me what to do' — emphasize advisory, "
        "the coach calls the play. (b) 'Compliance / privacy' — emphasize local-only, no uploads, no player PII.",
    )


def slide_close(prs: Presentation, total: int) -> None:
    s = add_blank(prs)
    fill_background(s, DARK)
    add_neuron_motif(s, int(Inches(1.0)), int(Inches(1.0)), scale=1.4)
    add_neuron_motif(s, int(Inches(11.0)), int(Inches(5.5)), scale=1.2)

    add_text(s, Inches(1.5), Inches(2.6), Inches(10), Inches(1.5),
             "Predict the swing\nbefore it happens.",
             size=58, color=ACCENT, bold=True, font=FONT_DISPLAY, align=PP_ALIGN.CENTER)

    add_text(s, Inches(1.5), Inches(4.6), Inches(10), Inches(0.5),
             "Ready to demo on a laptop today.",
             size=22, color=TEXT, align=PP_ALIGN.CENTER, italic=True)

    add_text(s, Inches(1.5), Inches(6.3), Inches(10), Inches(0.4),
             "Karsten Lowe  ·  kjlowe@ucsd.edu  ·  UCSD Triton Analytics",
             size=14, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

    add_footer(s, 7, total)
    add_notes(
        s,
        "End on the tagline. Offer the demo NOW — 'I have it open on my laptop right now if you want to "
        "see one rally.' That's the close. Don't leave the room without a concrete next step "
        "(pilot match date, intro to scout team, or follow-up meeting).",
    )


def build(out_path: Path) -> None:
    prs = make_prs()
    total = 7
    slide_title(prs, total)
    slide_problem(prs, total)
    slide_solution(prs, total)
    slide_credibility(prs, total)
    slide_workflow(prs, total)
    slide_ask(prs, total)
    slide_close(prs, total)
    prs.save(out_path)
    print(f"wrote {out_path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build(Path(__file__).parent / "FBSO_Predictor_Pitch.pptx")
